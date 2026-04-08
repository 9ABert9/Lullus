"""Document processor for multiple file formats.

Converts PDF, DOCX, PPTX, TXT, MD, HTML, EPUB, and CSV files into
a standardised Document dataclass for downstream processing.
"""

import csv
import logging
import re
import unicodedata
from dataclasses import dataclass, field
from datetime import datetime
from pathlib import Path
from typing import Any, Dict, List, Optional

logger = logging.getLogger(__name__)

SUPPORTED_FORMATS = {".pdf", ".docx", ".pptx", ".txt", ".md", ".html", ".htm", ".epub", ".csv"}


@dataclass
class Document:
    """Standardised representation of a processed document.

    Attributes:
        content: Full extracted text content.
        pages: List of per-page (or per-section) text strings.
        metadata: Dictionary with filename, format, page_count,
            file_size, date_processed, and any extra keys.
    """

    content: str
    pages: List[str]
    metadata: Dict[str, Any] = field(default_factory=dict)


class DocumentProcessingError(Exception):
    """Raised when a document cannot be processed."""


class DocumentProcessor:
    """Process documents of various formats into a unified Document object.

    Usage::

        processor = DocumentProcessor()
        doc = processor.process("/path/to/file.pdf")
        print(doc.content[:200])
    """

    # Map of file extension to processing method name
    _HANDLER_MAP: Dict[str, str] = {
        ".pdf": "_process_pdf",
        ".docx": "_process_docx",
        ".pptx": "_process_pptx",
        ".txt": "_process_txt",
        ".md": "_process_txt",  # Markdown is processed as plain text
        ".html": "_process_html",
        ".htm": "_process_html",
        ".epub": "_process_epub",
        ".csv": "_process_csv",
    }

    def process(self, file_path: str | Path) -> Document:
        """Process a file and return a Document.

        Args:
            file_path: Path to the file to process.

        Returns:
            A Document with extracted content, pages, and metadata.

        Raises:
            DocumentProcessingError: If the file cannot be processed.
            FileNotFoundError: If the file does not exist.
        """
        path = Path(file_path)
        if not path.exists():
            raise FileNotFoundError(f"File not found: {path}")

        ext = path.suffix.lower()
        if ext not in self._HANDLER_MAP:
            raise DocumentProcessingError(
                f"Unsupported format '{ext}'. Supported: {sorted(SUPPORTED_FORMATS)}"
            )

        handler_name = self._HANDLER_MAP[ext]
        handler = getattr(self, handler_name)

        logger.info("Processing '%s' (format=%s)", path.name, ext)
        try:
            pages: List[str] = handler(path)
        except DocumentProcessingError:
            raise
        except Exception as exc:
            logger.error("Error processing '%s': %s", path.name, exc)
            raise DocumentProcessingError(
                f"Failed to process '{path.name}': {exc}"
            ) from exc

        # Clean all pages
        pages = [self._clean_text(p) for p in pages]
        pages = [p for p in pages if p.strip()]

        content = "\n\n".join(pages)

        metadata: Dict[str, Any] = {
            "filename": path.name,
            "format": ext.lstrip("."),
            "page_count": len(pages),
            "file_size": path.stat().st_size,
            "date_processed": datetime.now().isoformat(),
        }

        logger.info(
            "Processed '%s': %d pages, %d chars",
            path.name,
            len(pages),
            len(content),
        )
        return Document(content=content, pages=pages, metadata=metadata)

    # ------------------------------------------------------------------
    # Format-specific processors
    # ------------------------------------------------------------------

    def _process_pdf(self, path: Path) -> List[str]:
        """Extract text from a PDF using PyMuPDF (fitz).

        Args:
            path: Path to the PDF file.

        Returns:
            List of strings, one per page.
        """
        try:
            import fitz  # PyMuPDF
        except ImportError as exc:
            raise DocumentProcessingError(
                "PyMuPDF is required for PDF processing. "
                "Install with: pip install PyMuPDF"
            ) from exc

        pages: List[str] = []
        try:
            doc = fitz.open(str(path))
            for page_num in range(len(doc)):
                page = doc.load_page(page_num)
                text = page.get_text("text")
                pages.append(text)
            doc.close()
        except Exception as exc:
            raise DocumentProcessingError(f"PDF extraction failed: {exc}") from exc

        logger.debug("PDF: extracted %d pages from '%s'", len(pages), path.name)
        return pages

    def _process_docx(self, path: Path) -> List[str]:
        """Extract text from a DOCX file using python-docx.

        Treats each paragraph as part of a single page. Section breaks
        are used to separate pages where possible.

        Args:
            path: Path to the DOCX file.

        Returns:
            List of strings (one per section or a single-element list).
        """
        try:
            from docx import Document as DocxDocument
        except ImportError as exc:
            raise DocumentProcessingError(
                "python-docx is required for DOCX processing. "
                "Install with: pip install python-docx"
            ) from exc

        try:
            doc = DocxDocument(str(path))
        except Exception as exc:
            raise DocumentProcessingError(f"DOCX extraction failed: {exc}") from exc

        # Collect paragraphs grouped by section breaks
        current_section_parts: List[str] = []
        sections: List[str] = []

        for para in doc.paragraphs:
            text = para.text.strip()
            if not text:
                continue
            # Detect section-break-like heading styles
            if para.style and para.style.name and para.style.name.startswith("Heading"):
                if current_section_parts:
                    sections.append("\n".join(current_section_parts))
                    current_section_parts = []
            current_section_parts.append(text)

        if current_section_parts:
            sections.append("\n".join(current_section_parts))

        if not sections:
            sections = [""]

        logger.debug("DOCX: extracted %d sections from '%s'", len(sections), path.name)
        return sections

    def _process_pptx(self, path: Path) -> List[str]:
        """Extract text from a PPTX file using python-pptx.

        Each slide becomes one page.

        Args:
            path: Path to the PPTX file.

        Returns:
            List of strings, one per slide.
        """
        try:
            from pptx import Presentation
        except ImportError as exc:
            raise DocumentProcessingError(
                "python-pptx is required for PPTX processing. "
                "Install with: pip install python-pptx"
            ) from exc

        try:
            prs = Presentation(str(path))
        except Exception as exc:
            raise DocumentProcessingError(f"PPTX extraction failed: {exc}") from exc

        pages: List[str] = []
        for slide_idx, slide in enumerate(prs.slides, start=1):
            slide_parts: List[str] = []
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for paragraph in shape.text_frame.paragraphs:
                        text = paragraph.text.strip()
                        if text:
                            slide_parts.append(text)
                if hasattr(shape, "table"):
                    table = shape.table
                    for row in table.rows:
                        row_text = " | ".join(
                            cell.text.strip() for cell in row.cells
                        )
                        if row_text.strip(" |"):
                            slide_parts.append(row_text)
            pages.append("\n".join(slide_parts))

        logger.debug("PPTX: extracted %d slides from '%s'", len(pages), path.name)
        return pages

    def _process_txt(self, path: Path) -> List[str]:
        """Read a plain text or Markdown file.

        Args:
            path: Path to the TXT/MD file.

        Returns:
            A single-element list with the file contents.
        """
        try:
            text = path.read_text(encoding="utf-8")
        except UnicodeDecodeError:
            try:
                text = path.read_text(encoding="latin-1")
                logger.warning("Fell back to latin-1 encoding for '%s'", path.name)
            except Exception as exc:
                raise DocumentProcessingError(
                    f"Could not decode text file: {exc}"
                ) from exc

        logger.debug("TXT/MD: read %d chars from '%s'", len(text), path.name)
        return [text]

    def _process_html(self, path: Path) -> List[str]:
        """Extract text from an HTML file using BeautifulSoup.

        Args:
            path: Path to the HTML file.

        Returns:
            A single-element list with the extracted text.
        """
        try:
            from bs4 import BeautifulSoup
        except ImportError as exc:
            raise DocumentProcessingError(
                "beautifulsoup4 is required for HTML processing. "
                "Install with: pip install beautifulsoup4"
            ) from exc

        try:
            raw = path.read_text(encoding="utf-8")
        except UnicodeDecodeError:
            raw = path.read_text(encoding="latin-1")
            logger.warning("Fell back to latin-1 encoding for '%s'", path.name)

        soup = BeautifulSoup(raw, "html.parser")

        # Remove script and style elements
        for element in soup(["script", "style", "nav", "footer", "header"]):
            element.decompose()

        text = soup.get_text(separator="\n")
        logger.debug("HTML: extracted %d chars from '%s'", len(text), path.name)
        return [text]

    def _process_epub(self, path: Path) -> List[str]:
        """Extract text from an EPUB file using ebooklib + BeautifulSoup.

        Each chapter/document becomes one page.

        Args:
            path: Path to the EPUB file.

        Returns:
            List of strings, one per chapter.
        """
        try:
            import ebooklib
            from ebooklib import epub
        except ImportError as exc:
            raise DocumentProcessingError(
                "ebooklib is required for EPUB processing. "
                "Install with: pip install ebooklib"
            ) from exc
        try:
            from bs4 import BeautifulSoup
        except ImportError as exc:
            raise DocumentProcessingError(
                "beautifulsoup4 is required for EPUB processing. "
                "Install with: pip install beautifulsoup4"
            ) from exc

        try:
            book = epub.read_epub(str(path), options={"ignore_ncx": True})
        except Exception as exc:
            raise DocumentProcessingError(f"EPUB extraction failed: {exc}") from exc

        pages: List[str] = []
        for item in book.get_items_of_type(ebooklib.ITEM_DOCUMENT):
            html_content = item.get_content().decode("utf-8", errors="replace")
            soup = BeautifulSoup(html_content, "html.parser")
            text = soup.get_text(separator="\n")
            if text.strip():
                pages.append(text)

        logger.debug("EPUB: extracted %d chapters from '%s'", len(pages), path.name)
        return pages

    def _process_csv(self, path: Path) -> List[str]:
        """Extract text from a CSV file using pandas.

        The entire CSV is converted to a readable string representation.
        Falls back to the csv stdlib module if pandas is unavailable.

        Args:
            path: Path to the CSV file.

        Returns:
            A single-element list with the CSV content as text.
        """
        try:
            import pandas as pd

            try:
                df = pd.read_csv(str(path), encoding="utf-8")
            except UnicodeDecodeError:
                df = pd.read_csv(str(path), encoding="latin-1")
                logger.warning("Fell back to latin-1 encoding for CSV '%s'", path.name)

            text = df.to_string(index=False)
            logger.debug("CSV (pandas): read %d rows from '%s'", len(df), path.name)
            return [text]

        except ImportError:
            logger.debug("pandas not available, falling back to csv module")

        # Fallback: stdlib csv
        rows: List[str] = []
        try:
            with open(path, "r", encoding="utf-8", newline="") as fh:
                reader = csv.reader(fh)
                for row in reader:
                    rows.append(" | ".join(row))
        except UnicodeDecodeError:
            with open(path, "r", encoding="latin-1", newline="") as fh:
                reader = csv.reader(fh)
                for row in reader:
                    rows.append(" | ".join(row))

        text = "\n".join(rows)
        logger.debug("CSV (stdlib): read %d rows from '%s'", len(rows), path.name)
        return [text]

    # ------------------------------------------------------------------
    # Text cleaning
    # ------------------------------------------------------------------

    @staticmethod
    def _clean_text(text: str) -> str:
        """Clean and normalize extracted text.

        Normalises unicode, collapses excessive whitespace, removes
        control characters, and strips leading/trailing whitespace.

        Args:
            text: Raw extracted text.

        Returns:
            Cleaned text.
        """
        # Normalize unicode (NFC form)
        text = unicodedata.normalize("NFC", text)

        # Remove control characters except newlines and tabs
        text = re.sub(r"[^\S\n\t]+", " ", text)

        # Collapse multiple blank lines into at most two newlines
        text = re.sub(r"\n{3,}", "\n\n", text)

        # Remove trailing whitespace on each line
        text = "\n".join(line.rstrip() for line in text.split("\n"))

        return text.strip()
